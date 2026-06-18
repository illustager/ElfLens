#!/usr/bin/env python3
"""Generate ElfLens project demonstration PPT (12 slides, dark theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Color palette (dark theme matching Avalonia Dark) ───────────────────
BG_DARK   = RGBColor(0x1E, 0x1E, 0x2E)
BG_CARD   = RGBColor(0x2A, 0x2A, 0x3E)
ACCENT    = RGBColor(0x4F, 0xC3, 0xF7)
WHITE     = RGBColor(0xE0, 0xE0, 0xE0)
GRAY      = RGBColor(0x90, 0x90, 0xAB)
GREEN     = RGBColor(0x81, 0xC7, 0x84)
RED       = RGBColor(0xEF, 0x53, 0x50)
ORANGE    = RGBColor(0xFF, 0xB7, 0x4D)
PURPLE    = RGBColor(0xCE, 0x93, 0xD8)
YELLOW    = RGBColor(0xFF, 0xE0, 0x82)
ADDR_GRAY = RGBColor(0x54, 0x6E, 0x7A)

FONT_BODY  = "Microsoft YaHei"
FONT_MONO  = "Cascadia Code"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── NSMAP for XML manipulation ─────────────────────────────────────────
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def set_text_margins(txBox, left=0.05, right=0.05, top=0.03, bottom=0.03):
    """Set internal margins on a text box to prevent text touching edges."""
    txBody = txBox._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if txBody is None:
        txBody = txBox._element.find('.//a:bodyPr', NSMAP)
    if txBody is not None:
        emu_per_inch = 914400
        txBody.set('lIns', str(int(left * emu_per_inch)))
        txBody.set('rIns', str(int(right * emu_per_inch)))
        txBody.set('tIns', str(int(top * emu_per_inch)))
        txBody.set('bIns', str(int(bottom * emu_per_inch)))

# ── Helper functions ────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height,
                text="", font_size=18, color=WHITE, bold=False,
                alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
                line_spacing=1.2, auto_fit=True):
    """Add a text box, optionally auto-sized. Returns (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    else:
        tf.auto_size = None

    set_text_margins(txBox)

    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = Pt(font_size * line_spacing)

    return txBox, tf


def add_para(tf, text, font_size=16, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
             space_after=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    return p


def add_bullet(tf, text, font_size=16, color=WHITE, bold=False,
               bullet_char="•", indent_level=0, space_after=4):
    """Add a bullet point. Use Unicode chars to avoid encoding issues."""
    p = add_para(tf, f"{bullet_char} {text}", font_size, color, bold,
                 space_after=space_after)
    p.level = indent_level
    return p


def add_card(slide, left, top, width, height, bg=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=12):
    add_textbox(slide, 11.8, 7.05, 1.2, 0.4,
                f"{num} / {total}", font_size=9, color=GRAY,
                alignment=PP_ALIGN.RIGHT, auto_fit=False)


def add_section_divider(slide, left, top, width, color=ACCENT):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_code_block(slide, left, top, width, height,
                   lines, font_size=10, line_spacing=1.5):
    """Add a code block with mono font. Card height is fixed, text box auto-grows."""
    card = add_card(slide, left, top, width, height)
    _, tf = add_textbox(slide, left + 0.15, top + 0.08,
                         width - 0.3, height - 0.16,
                         "", font_size=font_size, font_name=FONT_MONO,
                         color=WHITE, line_spacing=line_spacing,
                         auto_fit=False)
    # Remove default empty paragraph
    tf.paragraphs[0].text = ""
    first = True
    for line_info in lines:
        if isinstance(line_info, str):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = line_info
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.name = FONT_MONO
            p.space_after = Pt(0)
            p.space_before = Pt(0)
            first = False
        elif isinstance(line_info, list):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.font.size = Pt(font_size)
            p.font.name = FONT_MONO
            p.space_after = Pt(0)
            p.space_before = Pt(0)
            first = False
            for txt, clr in line_info:
                run = p.add_run()
                run.text = txt
                run.font.size = Pt(font_size)
                run.font.color.rgb = clr
                run.font.name = FONT_MONO
    return card


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_section_divider(s, 0, 0, 13.333, ACCENT)
s.shapes[-1].height = Pt(6)

add_textbox(s, 1.5, 1.8, 10.3, 1.2,
            "ElfLens", font_size=56, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 2.9, 10.3, 0.8,
            "跨平台远程 ELF 二进制调试器", font_size=28, color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 3.6, 10.3, 0.5,
            "A Cross-Platform Remote GDB Debugger Frontend",
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 4.6, 10.3, 0.5,
            "软件构造原理 (C#) · 课程期末项目  ·  v0.1.0",
            font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

tag_y = 5.3
tags = ["C#", ".NET 10", "Avalonia UI", "SSH.NET", "MVVM", "GPL v3.0"]
tag_colors = [GREEN, ACCENT, PURPLE, YELLOW, ORANGE, GRAY]
x = 1.8
for tag, tc in zip(tags, tag_colors):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(tag_y),
        Inches(1.5), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = tc
    shape.line.width = Pt(1)
    p = shape.text_frame.paragraphs[0]
    p.text = tag
    p.font.size = Pt(13)
    p.font.color.rgb = tc
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    x += 1.7

add_textbox(s, 1.5, 6.4, 10.3, 0.4,
            "illustager <illustager@outlook.com>  ·  2026.06",
            font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 1)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: PROJECT OVERVIEW
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "项目概述", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Left card: description
add_card(s, 0.8, 1.2, 5.5, 3.0)
_, tf = add_textbox(s, 1.0, 1.3, 5.1, 2.8,
                    "一句话描述", font_size=18, color=ACCENT, bold=True)
add_para(tf, "ElfLens 是一个跨平台 GDB 图形化调试前端，通过 SSH 连接远程 Linux 主机，"
         "为二进制分析提供一体化可视化工作区。", font_size=14, color=WHITE, space_after=12)
add_para(tf, "目标用户", font_size=18, color=ACCENT, bold=True, space_after=4)
for u in ["二进制逆向分析师", "CTF 竞赛选手", "漏洞安全研究者", "嵌入式系统开发者"]:
    add_bullet(tf, u, font_size=14, color=WHITE, space_after=3)

# Right card: key features
add_card(s, 6.8, 1.2, 5.7, 3.0)
_, tf = add_textbox(s, 7.0, 1.3, 5.3, 2.8,
                    "核心特色", font_size=18, color=ACCENT, bold=True)
features = [
    ("远程调试", "SSH 连接远程 Linux，无需本地二进制文件"),
    ("多面板工作区", "IDA Pro 风格的 4 区域停靠布局"),
    ("语法高亮", "自定义 Token 渲染管道，彩色反汇编代码"),
    ("端到端调试", "断点管理 -> 单步调试 -> 寄存器/栈分析全流程"),
    ("跨平台", "Windows / Linux / macOS 均可用"),
]
for title, desc in features:
    add_bullet(tf, f"{title}：{desc}", font_size=13, color=WHITE,
               bullet_char="▸", space_after=5)

# Bottom card: quick facts
add_card(s, 0.8, 4.5, 11.7, 2.7)
_, tf = add_textbox(s, 1.0, 4.6, 11.3, 2.5,
                    "项目速览", font_size=18, color=ACCENT, bold=True)
facts = [
    "项目类型：软件构造原理 (C#) 课程期末项目（单人开发）",
    "代码规模：约 30 个源文件，~4000 行 C# 代码，2 个 .NET 项目",
    "技术选型：.NET 10 + C# 12 + Avalonia UI + CommunityToolkit.Mvvm + SSH.NET",
    "开源协议：GNU General Public License v3.0",
    "开发周期：2026 年 6 月，Git 提交记录约 75 次",
]
for f in facts:
    add_bullet(tf, f, font_size=13, color=WHITE, space_after=5)

add_slide_number(s, 2)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: TECH STACK
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "技术栈", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Client-side tech stack table
add_card(s, 0.8, 1.2, 5.5, 5.7)
_, tf = add_textbox(s, 1.0, 1.3, 5.1, 0.5,
                    "客户端技术栈", font_size=20, color=ACCENT, bold=True)

table_data = [
    ("UI 框架", "Avalonia UI 12.0.4"),
    ("UI 主题", "Fluent Dark"),
    ("MVVM 框架", "CommunityToolkit.Mvvm 8.4.2"),
    ("SSH 库", "SSH.NET 2025.1.0"),
    ("运行时", ".NET 10.0"),
    ("语言", "C# (record, partial, GeneratedRegex)"),
    ("桌面运行时", "Avalonia.Desktop 12.0.4"),
    ("字体", "Cascadia Code / Inter"),
]

row_y = 2.1
for label, value in table_data:
    # label
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(row_y),
        Inches(2.0), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4E)
    shape.line.width = Pt(0.5)
    set_text_margins(shape, left=0.08, right=0.08, top=0.02, bottom=0.02)
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = PURPLE
    p.font.name = FONT_MONO
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT

    # value
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(row_y),
        Inches(3.1), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x4E)
    shape.line.width = Pt(0.5)
    set_text_margins(shape, left=0.08, right=0.08, top=0.02, bottom=0.02)
    p = shape.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(12)
    p.font.color.rgb = GREEN
    p.font.name = FONT_MONO
    p.alignment = PP_ALIGN.LEFT

    row_y += 0.46

# Remote host dependencies
add_card(s, 6.8, 1.2, 5.7, 2.4)
_, tf = add_textbox(s, 7.0, 1.3, 5.3, 2.2,
                    "远程主机依赖", font_size=20, color=ACCENT, bold=True)
remote_deps = [
    ("GDB", "GNU Debugger - 调试核心"),
    ("GNU Binutils", "objdump / readelf - 静态分析"),
    ("file", "文件类型识别"),
    ("pwndbg", "增强栈可视化 (可选)"),
    ("checksec", "安全属性检查 (可选)"),
]
for name, desc in remote_deps:
    add_bullet(tf, f"{name} — {desc}", font_size=14, color=WHITE, space_after=4)

# Project structure
add_card(s, 6.8, 3.9, 5.7, 3.0)
_, tf = add_textbox(s, 7.0, 4.0, 5.3, 2.8,
                    "项目结构", font_size=20, color=ACCENT, bold=True)
tree_lines = [
    ("ElfLens.slnx", GRAY),
    ("├── src/ElfLens/", GREEN),
    ("│   └── *.axaml + *.axaml.cs  (10 Views)", GRAY),
    ("├── src/ElfLens.Core/", ACCENT),
    ("│   ├── ViewModels/  (11 VMs)", GRAY),
    ("│   ├── Services/    (2 services)", GRAY),
    ("│   └── Models/      (2 models)", GRAY),
    ("└── docs/  ·  tests/  ·  release/", GRAY),
]
for text, color in tree_lines:
    add_para(tf, text, font_size=13, color=color, font_name=FONT_MONO, space_after=3)

add_slide_number(s, 3)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "系统架构", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Layered architecture
add_card(s, 0.8, 1.2, 7.0, 5.8)
_, tf = add_textbox(s, 1.0, 1.3, 6.4, 0.5,
                    "分层架构 (Layered Architecture)", font_size=20, color=ACCENT, bold=True)

layers = [
    ("Presentation Layer  (ElfLens)", GREEN,
     ["Avalonia Views (.axaml) · Data Binding",
      "Converters · ViewLocator · Dark Theme",
      "10 个用户控件视图 + 1 个主窗口"]),
    ("ViewModel Layer  (ElfLens.Core)", ACCENT,
     ["MainViewModel - 应用编排中心",
      "PanelViewModel 体系 - 7 个可停靠面板",
      "CommunityToolkit.Mvvm 源码生成器",
      "[ObservableProperty] + [RelayCommand]"]),
    ("Service Layer  (ElfLens.Core)", PURPLE,
     ["ISshService 接口 -> SshService 实现",
      "SSH.NET 封装：连接 / 命令 / Shell 流",
      "异步模式 (Task.Run + Async/Await)"]),
    ("Model Layer  (ElfLens.Core)", YELLOW,
     ["ShellSession - SSH Shell 流包装器",
      "SshConnectionInfo - 连接参数 DTO",
      "Display Types: Token, HighlightedLine, FunctionItem"]),
]

layer_y = 2.1
for name, color, items in layers:
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(layer_y), Inches(6.4), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    set_text_margins(shape, left=0.08, right=0.08, top=0.02, bottom=0.02)
    p = shape.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(12)
    p.font.color.rgb = BG_DARK if color == YELLOW else WHITE
    p.font.name = FONT_MONO
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    layer_y += 0.48
    for item in items:
        add_textbox(s, 1.3, layer_y, 5.8, 0.28,
                    f"  {item}", font_size=11, color=GRAY, auto_fit=False)
        layer_y += 0.28
    layer_y += 0.12

# Right: MVVM data flow
add_card(s, 8.3, 1.2, 4.4, 2.8)
_, tf = add_textbox(s, 8.5, 1.3, 4.0, 2.6,
                    "MVVM 数据流", font_size=18, color=ACCENT, bold=True)
flow_items = [
    ("View (.axaml)", GREEN),
    ("  ↕  Data Binding (Compiled)", GRAY),
    ("ViewModel (ElfLens.Core)", ACCENT),
    ("  ↕  Method Calls", GRAY),
    ("Model / Service", YELLOW),
    ("", GRAY),
    ("ViewLocator: XxxViewModel -> XxxView", GRAY),
]
for text, color in flow_items:
    if text:
        add_para(tf, text, font_size=13 if text.strip().startswith("↕") or "Xxx" in text else 14,
                 color=color, bold=(color not in (GRAY,)), font_name=FONT_MONO, space_after=5)
    else:
        add_para(tf, "", font_size=6, space_after=2)

# Right: Event system
add_card(s, 8.3, 4.3, 4.4, 2.7)
_, tf = add_textbox(s, 8.5, 4.4, 4.0, 2.5,
                    "事件驱动协调", font_size=18, color=ACCENT, bold=True)
events = [
    ("SessionChanged", "GDB 会话 -> 所有依赖面板"),
    ("BreakpointRequested", "右键菜单 -> 断点面板"),
    ("BlocksChanged", "新函数块 -> 重新标记断点"),
    ("OnChanged", "断点变更 -> 反汇编视图更新"),
    ("NavigateTo", "点击引用 -> 滚动定位"),
]
for evt, desc in events:
    add_para(tf, evt, font_size=12, color=ORANGE, bold=True,
             font_name=FONT_MONO, space_after=1)
    add_para(tf, f"  {desc}", font_size=11, color=GRAY, space_after=5)

add_slide_number(s, 4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: CORE MODULES
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "核心功能模块", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

modules = [
    ("SSH 远程连接",
     "支持密码 / 私钥认证\nSSH.NET 封装\n异步连接与 Shell 流管理",
     GREEN),
    ("ELF 文件分析",
     "file / readelf / checksec\nELF 头、段表、节表\n安全属性 (NX/RELRO/PIE/Canary)",
     PURPLE),
    ("静态反汇编",
     "objdump -d 解析\n折叠式函数视图\n点击导航函数引用",
     ACCENT),
    ("交互式 GDB 调试",
     "Step In / Step Over / Continue\n实时反汇编 (disassemble /r)\n当前 PC 自动高亮与展开",
     RED),
    ("断点管理",
     "按函数名 / 地址添加断点\n启用 / 禁用 / 删除 / 批量应用\n反汇编视图断点标记",
     ORANGE),
    ("寄存器 & 调用栈",
     "info registers 寄存器展示\nbt 回溯 + info frame\npwndbg/x/gx 栈内存转储",
     YELLOW),
]

for i, (title, desc, color) in enumerate(modules):
    col = i % 3
    row = i // 3
    left = 0.8 + col * 4.15
    top = 1.2 + row * 3.05

    add_card(s, left, top, 3.85, 2.85)
    _, tf = add_textbox(s, left + 0.15, top + 0.1, 3.55, 2.65,
                        title, font_size=18, color=color, bold=True)
    for line in desc.split("\n"):
        add_bullet(tf, line, font_size=13, color=WHITE, space_after=4)

add_slide_number(s, 5)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: WORKSPACE LAYOUT
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "工作区布局", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

DARK_INNER = RGBColor(0x1A, 0x1A, 0x2A)

# Left zone
add_card(s, 0.8, 1.4, 2.8, 3.6, BG_CARD)
_, tf = add_textbox(s, 0.95, 1.45, 2.5, 0.35,
                    "Left  (220px)", font_size=12, color=GRAY, bold=True, auto_fit=False)
# Registers
add_card(s, 0.95, 1.95, 2.5, 1.35, DARK_INNER)
_, tf = add_textbox(s, 1.05, 2.15, 2.3, 0.9,
                    "Registers\n寄存器面板", font_size=13, color=GREEN, alignment=PP_ALIGN.CENTER)
# Stack
add_card(s, 0.95, 3.5, 2.5, 1.35, DARK_INNER)
_, tf = add_textbox(s, 1.05, 3.7, 2.3, 0.9,
                    "Stack\n调用栈面板", font_size=13, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Center zone
add_card(s, 3.9, 1.4, 5.5, 3.6, BG_CARD)
_, tf = add_textbox(s, 4.05, 1.45, 5.2, 0.35,
                    "Center (flexible)", font_size=12, color=GRAY, bold=True, auto_fit=False)
# Disassembly
add_card(s, 4.05, 1.95, 5.2, 1.35, DARK_INNER)
_, tf = add_textbox(s, 4.15, 2.2, 5.0, 0.8,
                    "Disassembly - 静态反汇编 (objdump)", font_size=12, color=ACCENT, alignment=PP_ALIGN.CENTER)
# GDB
add_card(s, 4.05, 3.5, 5.2, 1.35, DARK_INNER)
_, tf = add_textbox(s, 4.15, 3.75, 5.0, 0.8,
                    "GDB Disasm - 交互式调试 (gdb disassemble /r)", font_size=12, color=RED, alignment=PP_ALIGN.CENTER)

# Right zone
add_card(s, 9.6, 1.4, 2.9, 3.6, BG_CARD)
_, tf = add_textbox(s, 9.75, 1.45, 2.6, 0.35,
                    "Right  (220px)", font_size=12, color=GRAY, bold=True, auto_fit=False)
# File Info
add_card(s, 9.75, 1.95, 2.6, 1.35, DARK_INNER)
_, tf = add_textbox(s, 9.85, 2.2, 2.4, 0.8,
                    "File Info\nELF 文件信息", font_size=12, color=PURPLE, alignment=PP_ALIGN.CENTER)
# Breakpoints
add_card(s, 9.75, 3.5, 2.6, 1.35, DARK_INNER)
_, tf = add_textbox(s, 9.85, 3.75, 2.4, 0.8,
                    "Breakpoints\n断点管理", font_size=12, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Bottom zone
add_card(s, 0.8, 5.3, 11.7, 1.95, BG_CARD)
_, tf = add_textbox(s, 0.95, 5.35, 11.4, 0.35,
                    "Bottom  (200px)", font_size=12, color=GRAY, bold=True, auto_fit=False)

half_w = 5.5
add_card(s, 0.95, 5.85, half_w, 1.2, DARK_INNER)
_, tf = add_textbox(s, 1.05, 6.05, half_w - 0.2, 0.8,
                    "SSH Shell\n通用 SSH 终端", font_size=12, color=GREEN, alignment=PP_ALIGN.CENTER)

add_card(s, 1.15 + half_w, 5.85, half_w, 1.2, DARK_INNER)
_, tf = add_textbox(s, 1.25 + half_w, 6.05, half_w - 0.2, 0.8,
                    "GDB Shell\nGDB 交互终端 (自动创建)", font_size=12, color=RED, alignment=PP_ALIGN.CENTER)

add_slide_number(s, 6)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: DISASSEMBLY & SYNTAX HIGHLIGHTING
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6,
            "反汇编与语法高亮", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Rendering pipeline
add_card(s, 0.8, 1.2, 7.0, 2.6)
_, tf = add_textbox(s, 1.0, 1.3, 6.6, 2.4,
                    "自定义 Token 渲染管道", font_size=20, color=ACCENT, bold=True)
add_para(tf, "原始工具输出 -> Tokenize(正则解析) -> List<Token> -> ItemsControl 渲染",
         font_size=14, color=WHITE, font_name=FONT_MONO)
add_para(tf, "", font_size=4, space_after=2)
pipeline_steps = [
    "1. objdump / GDB 输出原始文本行",
    "2. DisassemblyHighlighter.Tokenize() 正则解析为 Token 列表",
    "3. 每个 Token = (Text, Color, NavigateTo?)",
    "4. UI 通过 ItemsControl 水平排列彩色 TextBlock",
    "5. 点击可导航 Token 触发 NavigateToFunction 事件",
]
for step in pipeline_steps:
    add_bullet(tf, step, font_size=13, color=GRAY, space_after=3)

# Color scheme table
add_card(s, 8.3, 1.2, 4.4, 2.6)
_, tf = add_textbox(s, 8.5, 1.3, 4.0, 0.5,
                    "配色方案", font_size=18, color=ACCENT, bold=True)

color_table = [
    ("address", "546E7A", ADDR_GRAY),
    ("call", "81C784", GREEN),
    ("ret", "EF5350", RED),
    ("branch", "FFB74D", ORANGE),
    ("register", "CE93D8", PURPLE),
    ("hex value", "FFE082", YELLOW),
    ("func ref", "4FC3F7", ACCENT),
    ("comment", "6A9955", RGBColor(0x6A, 0x99, 0x55)),
]

row_y = 2.0
n_rows = len(color_table)
# Ensure rows fit in card: card bottom = 1.2 + 2.6 = 3.8, but we need space
# row_y starts at 2.0, each row is 0.29", 8 rows = 2.32", end = 4.32 > 3.8
# Fix: reduce row height and font
row_h = 0.25
row_font = 10

for label, hex_val, clr in color_table:
    # label
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(row_y),
        Inches(1.0), Inches(row_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.fill.background()
    set_text_margins(shape, left=0.04, right=0.04, top=0.01, bottom=0.01)
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(row_font)
    p.font.color.rgb = GRAY
    p.font.name = FONT_MONO
    p.alignment = PP_ALIGN.LEFT

    # swatch
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.65), Inches(row_y + 0.02),
        Inches(0.22), Inches(row_h - 0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()

    # hex
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(10.0), Inches(row_y),
        Inches(1.2), Inches(row_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.fill.background()
    set_text_margins(shape, left=0.04, right=0.04, top=0.01, bottom=0.01)
    p = shape.text_frame.paragraphs[0]
    p.text = f"#{hex_val}"
    p.font.size = Pt(row_font)
    p.font.color.rgb = clr
    p.font.name = FONT_MONO
    p.alignment = PP_ALIGN.LEFT

    row_y += 0.29

# Code example
code_lines = [
    "  4048b0 <main>:",
    [("  4048b0:", ADDR_GRAY), (" 55", ADDR_GRAY),
     ("     ", ADDR_GRAY), ("push", WHITE),
     ("   ", ADDR_GRAY), ("%rbp", PURPLE)],
    [("  4048b1:", ADDR_GRAY), (" 48 89 e5", ADDR_GRAY),
     ("  ", ADDR_GRAY), ("mov", WHITE),
     ("    ", ADDR_GRAY), ("%rsp", PURPLE), (",", WHITE), ("%rbp", PURPLE)],
    [("  4048b4:", ADDR_GRAY), (" 48 83 ec 10", ADDR_GRAY),
     ("  ", ADDR_GRAY), ("sub", WHITE),
     ("    ", ADDR_GRAY), ("$0x10", YELLOW), (",", WHITE), ("%rsp", PURPLE)],
    [("  4048b8:", ADDR_GRAY), (" e8 a3 ff ff ff", ADDR_GRAY),
     ("  ", ADDR_GRAY), ("call", GREEN),
     ("   ", ADDR_GRAY), ("<init>", ACCENT)],
    [("  4048bd:", ADDR_GRAY), (" b8 00 00 00 00", ADDR_GRAY),
     ("  ", ADDR_GRAY), ("mov", WHITE),
     ("    ", ADDR_GRAY), ("$0x0", YELLOW), (",", WHITE), ("%eax", PURPLE)],
    [("  4048c2:", ADDR_GRAY), (" c3", ADDR_GRAY),
     ("     ", ADDR_GRAY), ("ret", RED)],
]
add_code_block(s, 0.8, 4.1, 11.7, 3.1, code_lines, font_size=12)

add_slide_number(s, 7)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: GDB INTERACTIVE DEBUGGING
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6,
            "交互式 GDB 调试", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Debug lifecycle
add_card(s, 0.8, 1.2, 5.8, 2.6)
_, tf = add_textbox(s, 1.0, 1.3, 5.4, 2.4,
                    "调试生命周期", font_size=20, color=ACCENT, bold=True)
flow_steps = [
    "1. 用户点击 Start Debugging",
    "2. 通过 ShellSession 启动 gdb -q <binary>",
    "3. 发送 run 命令启动目标程序",
    "4. 触发 SessionChanged -> MainViewModel 连接所有面板",
    "5. 批量应用已有断点",
    "6. 读取 info registers pc 获取当前 PC",
    "7. disassemble /r 获取实时反汇编",
    "8. 高亮当前指令 -> 自动展开所在函数 -> 滚动定位",
]
for step in flow_steps:
    add_bullet(tf, step, font_size=12, color=GRAY, space_after=3)

# Debug controls
add_card(s, 7.1, 1.2, 5.4, 2.6)
_, tf = add_textbox(s, 7.3, 1.3, 5.0, 2.4,
                    "调试控制", font_size=20, color=ACCENT, bold=True)
controls = [
    ("Step Into  (stepi)", "逐指令进入函数调用"),
    ("Step Over  (nexti)", "逐指令跳过函数调用"),
    ("Continue   (continue)", "继续执行至下一断点"),
    ("Restart", "重新运行程序并恢复所有断点"),
    ("Stop", "中断执行"),
    ("Set BP (右键)", "在反汇编视图中右键设置断点"),
]
for cmd, desc in controls:
    add_para(tf, f"▸ {cmd}", font_size=12, color=GREEN, bold=True,
             font_name=FONT_MONO, space_after=1)
    add_para(tf, f"   {desc}", font_size=11, color=GRAY, space_after=5)

# Function blocks
add_card(s, 0.8, 4.1, 5.8, 3.1)
_, tf = add_textbox(s, 1.0, 4.2, 5.4, 2.9,
                    "函数块管理", font_size=18, color=ACCENT, bold=True)
block_items = [
    "FunctionItem 模型支持折叠/展开",
    "首次访问函数时懒加载反汇编 (lazy fetch)",
    "已访问函数缓存于 FunctionBlocks 集合中",
    "当前执行函数自动展开",
    "PC 所在行黄色背景高亮 (CurrentBg 转换器)",
    "断点行红色/橙色左边框标记",
    "ScrollToBlock 事件驱动自动滚动定位",
]
for item in block_items:
    add_bullet(tf, item, font_size=13, color=GRAY, space_after=4)

# GDB disassembly format
add_card(s, 7.1, 4.1, 5.4, 3.1)
gdb_code = [
    "GDB disassemble /r 输出示例：",
    "",
    [("   0x4048b0 <main+0>", ADDR_GRAY), (" 55", ADDR_GRAY),
     ("\t", ADDR_GRAY), ("push", WHITE), ("   %rbp", PURPLE)],
    [("   0x4048b1 <main+1>", ADDR_GRAY), (" 48 89 e5", ADDR_GRAY),
     ("\t", ADDR_GRAY), ("mov", WHITE), ("    %rsp,%rbp", PURPLE)],
    [("=> 0x4048b4 <main+4>", YELLOW), (" 48 83 ec 10", ADDR_GRAY),
     ("\t", ADDR_GRAY), ("sub", WHITE), ("    $0x10,%rsp", PURPLE)],
    [("   0x4048b8 <main+8>", ADDR_GRAY), (" e8 a3 ff ff ff", ADDR_GRAY),
     ("\t", ADDR_GRAY), ("call", WHITE), ("  <init>", ACCENT)],
    [("   0x4048bd <main+13>", ADDR_GRAY), (" b8 00 00 00 00", ADDR_GRAY),
     ("\t", ADDR_GRAY), ("mov", WHITE), ("   $0x0,%eax", PURPLE)],
    [("   0x4048c2 <main+18>", ADDR_GRAY), (" c3", ADDR_GRAY),
     ("\t", ADDR_GRAY), ("ret", WHITE)],
    "",
    "=> 黄色 = 当前 PC 位置",
]
add_code_block(s, 7.1, 4.1, 5.4, 3.1, gdb_code, font_size=10)

add_slide_number(s, 8)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: SSH COMMUNICATION DESIGN
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6,
            "SSH 远程通信设计", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# ShellSession design
add_card(s, 0.8, 1.2, 6.3, 3.6)
_, tf = add_textbox(s, 1.0, 1.3, 5.9, 3.4,
                    "ShellSession - 核心抽象", font_size=20, color=ACCENT, bold=True)

add_para(tf, "设计问题：GDB 是有状态交互式程序，无法用简单的 SSH 单次命令执行",
         font_size=12, color=GRAY, space_after=10)
add_para(tf, "ShellSession 解决方案：", font_size=14, color=WHITE, bold=True, space_after=6)

design_points = [
    "封装 SSH.NET ShellStream -> xterm-256color 伪终端",
    "后台读取循环：持续从 SSH 流读取输出",
    "ANSI 转义序列过滤：[GeneratedRegex] 编译时生成正则",
    "CaptureOutputAsync() 统一命令/响应模式",
    "  发送命令 -> 累积输出 -> 停止条件匹配 -> 返回",
    "停止条件：GDB 提示符 | 自定义谓词 | 8000 字符 | 800ms 超时",
    "SemaphoreSlim 写锁：序列化命令发送，防止交错",
    "OnOutput 事件：实时输出通知（Shell 面板渲染用）",
]
for point in design_points:
    add_bullet(tf, point, font_size=12, color=GRAY, space_after=3)

# Three access modes
add_card(s, 7.6, 1.2, 4.9, 3.6)
_, tf = add_textbox(s, 7.8, 1.3, 4.5, 3.4,
                    "三种访问模式", font_size=18, color=ACCENT, bold=True)

modes = [
    ("ExecuteCommandAsync",
     ["一次性命令执行",
      "file / readelf / objdump",
      "SshClient.CreateCommand()",
      "30 秒超时"]),
    ("CreateShellSessionAsync",
     ["持久化 Shell 流",
      "通用 SSH 终端面板",
      "命令历史 / 实时交互"]),
    ("GDB Shell Reuse",
     ["复用已有 ShellSession",
      "GDB Shell 面板 (调试时自动创建)",
      "共享 GDB 进程上下文"]),
]
for name, lines in modes:
    add_para(tf, f"▸ {name}", font_size=13, color=GREEN, bold=True,
             font_name=FONT_MONO, space_after=1)
    for line in lines:
        add_para(tf, f"   {line}", font_size=11, color=GRAY, space_after=1)
    add_para(tf, "", font_size=4, space_after=4)

# Code sample
code_lines = [
    "// CaptureOutputAsync - 所有面板的统一命令接口",
    "",
    [("public async ", WHITE), ("Task", RGBColor(0x56, 0x9C, 0xD6)),
     ("<", WHITE), ("string", RGBColor(0x56, 0x9C, 0xD6)),
     ("> ", WHITE),
     ("CaptureOutputAsync", YELLOW), ("(", WHITE)],
    [("    ", WHITE), ("string", RGBColor(0x56, 0x9C, 0xD6)),
     (" command,", WHITE)],
    [("    ", WHITE), ("int", RGBColor(0x56, 0x9C, 0xD6)),
     (" timeoutMs = 800,", WHITE)],
    [("    ", WHITE), ("Func", RGBColor(0x56, 0x9C, 0xD6)),
     ("<string, bool> stopPredicate = null)", WHITE)],
    "",
    "// 使用示例：向 GDB 发送命令并等待响应",
    "",
    [("var output = await session.CaptureOutputAsync(", WHITE)],
    [("    ", WHITE), ('"info breakpoints"', RGBColor(0xD6, 0x9D, 0x85)), (",", WHITE)],
    [("    stopPredicate: s => s.Contains(", WHITE),
     ('"No breakpoints"', RGBColor(0xD6, 0x9D, 0x85)), ("));", WHITE)],
]
add_code_block(s, 0.8, 5.1, 11.7, 2.1, code_lines, font_size=11)

add_slide_number(s, 9)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: ENGINEERING HIGHLIGHTS
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "工程亮点", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Redesigned: 4 highlights as 2x2 grid with generous card sizes
highlights = [
    ("MVVM 源码生成器", ACCENT,
     ["[ObservableProperty] 自动生成属性变更通知代码",
      "[RelayCommand] 自动生成 ICommand 实现",
      "[GeneratedRegex] 编译时生成 ANSI 过滤正则",
      "消除大量样板代码，提升可维护性"],
     ["partial class ConnectPageViewModel",
      "{",
      "    [ObservableProperty]",
      '    private string _host = "";',
      "    // auto-gen: public string Host { get; set; }",
      "    //           + PropertyChanged 通知",
      "}"]),
    ("事件驱动面板协调", GREEN,
     ["MainViewModel 作为中心编排器",
      "7 个面板通过事件解耦通信",
      "SessionChanged 一键连接所有 GDB 依赖面板",
      "断点变更自动刷新两个反汇编视图"],
     ["// GDB session -> all dependent panels",
      "gdbDisasm.SessionChanged += (s, e) => {",
      "    breakpointPanel.SetSession(e.Session);",
      "    registersPanel.SetSession(e.Session);",
      "    stackPanel.SetSession(e.Session);",
      "    AddOrRemoveGdbShellTab(e);",
      "};"]),
    ("Token 导航系统", PURPLE,
     ["可点击的 Token -> NavigateTo 属性",
      "函数引用 <func_name> 点击直接跳转",
      "跨面板联动：断点面板 -> 反汇编视图标记",
      "自定义 IValueConverter 处理光标/颜色/边框"],
     ["// Token 记录类型",
      "record Token(string Text, string Color,",
      "            string? NavigateTo);",
      "",
      "// 渲染: StackPanel 水平排列 TextBlock",
      "// 点击: NavigateToFunction 事件 -> 滚动定位"]),
    ("渐进增强策略", ORANGE,
     ["checksec 不可用 -> 手动解析 readelf -lW/-sW",
      "pwndbg 不可用 -> 回退到原生 GDB x/gx 命令",
      "StackHighlighter 自动检测两种输出格式",
      "双模式 ShellPanel：SSH 初始化 vs GDB 复用"],
     ["// StackHighlighter 双格式支持",
      "if (line.Contains('|') || line.Contains('<'))",
      "    return ParsePwndbg(line); // pwndbg",
      "else",
      "    return ParseGdbHex(line);  // x/gx fallback"]),
]

for i, (title, color, desc_lines, code) in enumerate(highlights):
    col = i % 2
    row = i // 2
    left = 0.8 + col * 6.35
    top = 1.15 + row * 3.15
    card_w = 5.95
    card_h = 2.95

    add_card(s, left, top, card_w, card_h)

    # Title
    _, tf = add_textbox(s, left + 0.15, top + 0.1, card_w - 0.3, 0.35,
                        title, font_size=17, color=color, bold=True)

    # Description bullets (top portion)
    _, tf_desc = add_textbox(s, left + 0.15, top + 0.55, card_w - 0.3, 1.0,
                             "", font_size=10, color=GRAY)
    for line in desc_lines:
        add_bullet(tf_desc, line, font_size=10, color=GRAY, space_after=2,
                   bullet_char="•")

    # Code snippet (bottom portion)
    code_card = add_card(s, left + 0.15, top + 1.5, card_w - 0.3, card_h - 1.65,
                         RGBColor(0x1A, 0x1A, 0x2A))
    _, tf_code = add_textbox(s, left + 0.25, top + 1.55,
                              card_w - 0.5, card_h - 1.75,
                              "", font_size=9, font_name=FONT_MONO,
                              color=WHITE, line_spacing=1.4, auto_fit=False)
    tf_code.paragraphs[0].text = ""
    first = True
    for line_info in code:
        if isinstance(line_info, str):
            p = tf_code.paragraphs[0] if first else tf_code.add_paragraph()
            p.text = line_info
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY
            p.font.name = FONT_MONO
            p.space_after = Pt(0)
            p.space_before = Pt(0)
            first = False
        elif isinstance(line_info, list):
            p = tf_code.paragraphs[0] if first else tf_code.add_paragraph()
            p.font.size = Pt(9)
            p.font.name = FONT_MONO
            p.space_after = Pt(0)
            p.space_before = Pt(0)
            first = False
            for txt, clr in line_info:
                run = p.add_run()
                run.text = txt
                run.font.size = Pt(9)
                run.font.color.rgb = clr
                run.font.name = FONT_MONO

add_slide_number(s, 10)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: PROJECT SUMMARY
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_textbox(s, 0.8, 0.3, 11.7, 0.6, "项目总结", font_size=32, color=ACCENT, bold=True)
add_section_divider(s, 0.8, 0.9, 2.5, ACCENT)

# Completed features
add_card(s, 0.8, 1.2, 5.8, 3.8)
_, tf = add_textbox(s, 1.0, 1.3, 5.4, 3.6,
                    "已完成功能", font_size=20, color=GREEN, bold=True)
done = [
    "SSH 远程连接 (密码 + 私钥认证)",
    "ELF 文件分析 (file / readelf / checksec / 手动回退)",
    "静态反汇编 (objdump 解析 + 函数折叠)",
    "汇编语法高亮 (8 种 Token 颜色，可点击导航)",
    "交互式 GDB 调试 (Step In/Over, Continue, Restart, Stop)",
    "断点管理 (函数名/地址添加，启用/禁用/删除)",
    "断点可视化标记 (反汇编视图红色/橙色左边框)",
    "寄存器面板 (info registers 解析)",
    "调用栈面板 (bt + 可展开帧 + 栈内存转储)",
    "栈内存语法高亮 (pwndbg + GDB x/gx 双格式)",
    "交互式终端 (SSH Shell + GDB Shell，命令历史)",
    "事件驱动面板协调 (7 面板 + 7 种事件)",
    "跨平台支持 (Windows / Linux / macOS)",
]
for item in done:
    add_bullet(tf, item, font_size=12, color=WHITE, space_after=3,
               bullet_char="✓")

# Project stats
add_card(s, 7.1, 1.2, 5.4, 3.8)
_, tf = add_textbox(s, 7.3, 1.3, 5.0, 3.6,
                    "项目统计", font_size=20, color=ACCENT, bold=True)

stats = [
    ("C# 源文件", "~30 个"),
    ("代码行数", "~4,000 行"),
    (".NET 项目", "2 个"),
    ("Views", "10 个 .axaml"),
    ("ViewModels", "11 个"),
    ("Models", "2 个"),
    ("Services", "2 个"),
    ("Highlighters", "2 个"),
    ("NuGet 依赖", "5 个"),
    ("Git 提交", "~75 次"),
    ("测试覆盖", "暂无"),
]

# Use a 2-column mini-table inside the text box
# Simpler: list stats as label: value pairs
stats_text = [
    ("C# 源文件：约 30 个", "代码行数：约 4,000 行"),
    (".NET 项目：2 个", "Views：10 个 .axaml"),
    ("ViewModels：11 个", "Models + Services：4 个"),
    ("Highlighters：2 个", "NuGet 直接依赖：7 个"),
    ("Git 提交：约 75 次", "测试覆盖：暂无"),
]
for left_line, right_line in stats_text:
    add_para(tf, f"  {left_line}     {right_line}", font_size=12,
             color=WHITE, font_name=FONT_MONO, space_after=6)

add_para(tf, "", font_size=8, space_after=4)
add_para(tf, "架构：MVVM + 分层架构 + 事件驱动", font_size=12,
         color=GRAY, space_after=3)
add_para(tf, "技术：C# 12 / .NET 10 / Avalonia UI 12", font_size=12,
         color=GRAY, space_after=3)
add_para(tf, "许可：GNU General Public License v3.0", font_size=12,
         color=PURPLE, space_after=3)

# Learnings
add_card(s, 0.8, 5.25, 5.8, 2.0)
_, tf = add_textbox(s, 1.0, 5.35, 5.4, 1.8,
                    "技术收获", font_size=18, color=ACCENT, bold=True)
learnings = [
    "MVVM 架构设计与 CommunityToolkit.Mvvm 源码生成器的深度使用",
    "Avalonia UI 跨平台桌面应用开发（编译绑定、DataTemplate）",
    "SSH.NET 异步封装与 Shell 流的状态管理",
    "编译时正则生成 ([GeneratedRegex]) 提升性能",
    "事件驱动架构解耦复杂面板交互",
]
for item in learnings:
    add_bullet(tf, item, font_size=11, color=GRAY, space_after=3)

# Future
add_card(s, 7.1, 5.25, 5.4, 2.0)
_, tf = add_textbox(s, 7.3, 5.35, 5.0, 1.8,
                    "后续方向", font_size=18, color=ACCENT, bold=True)
future = [
    "编写单元测试覆盖核心逻辑",
    "支持 GDB/MI 替代 Shell 解析",
    "添加内存查看器与 Watches 面板",
    "支持更多架构 (ARM / RISC-V)",
    "本地调试模式 (无需 SSH)",
    "插件系统支持自定义面板",
]
for item in future:
    add_bullet(tf, item, font_size=11, color=GRAY, space_after=3,
               bullet_char="→")

add_slide_number(s, 11)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12: THANKS / Q&A
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)

add_section_divider(s, 0, 0, 13.333, ACCENT)
s.shapes[-1].height = Pt(6)

add_textbox(s, 1.5, 1.8, 10.3, 1.0,
            "感谢聆听", font_size=52, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 2.8, 10.3, 0.6,
            "Questions & Answers", font_size=24, color=GRAY,
            alignment=PP_ALIGN.CENTER)

add_card(s, 3.5, 3.9, 6.3, 2.0)
_, tf = add_textbox(s, 3.7, 4.0, 5.9, 1.8,
                    "项目信息", font_size=20, color=ACCENT, bold=True,
                    alignment=PP_ALIGN.CENTER)
add_para(tf, "ElfLens v0.1.0", font_size=17, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER, space_after=8)
add_para(tf, "软件构造原理 (C#) · 课程期末项目", font_size=14, color=GRAY,
         alignment=PP_ALIGN.CENTER, space_after=8)
add_para(tf, "开源协议：GNU General Public License v3.0", font_size=13, color=PURPLE,
         alignment=PP_ALIGN.CENTER, space_after=8)
add_para(tf, "Author: illustager <illustager@outlook.com>", font_size=13, color=GREEN,
         alignment=PP_ALIGN.CENTER, space_after=8)
add_para(tf, "Tech: C# 12 + .NET 10 + Avalonia UI + SSH.NET", font_size=12, color=GRAY,
         alignment=PP_ALIGN.CENTER)

add_slide_number(s, 12)

# ── Save ────────────────────────────────────────────────────────────────
output_path = "ElfLens-demo.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
